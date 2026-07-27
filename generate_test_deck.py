"""
Test Slide Generation Script for PowerScale PPTX Architecture
"""

import os
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from powerscale_pptx import (
    create_presentation, set_slide_background, add_header, add_card,
    add_kpi_block, add_footer, calculate_grid_positions,
    BLUE, CARD_BG, CARD_BORDER, TEXT_BODY, TEXT_MUTED, CANVAS
)

def build_test_deck(output_filename="powerscale_demo.pptx"):
    prs = create_presentation()
    blank_layout = prs.slide_layouts[6] # Blank layout

    # ==========================================
    # Slide 1: Title Slide
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, CANVAS)

    # Title Card
    card_w = Inches(11.733)
    card_h = Inches(5.2)
    card_l = Inches(0.8)
    card_t = Inches(1.1)

    add_card(slide1, card_l, card_t, card_w, card_h, bg_color=CARD_BG, border_color=CARD_BORDER)

    # Text container inside Card
    txBox = slide1.shapes.add_textbox(card_l + Inches(0.8), card_t + Inches(1.0), card_w - Inches(1.6), card_h - Inches(2.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)

    p_eye = tf.paragraphs[0]
    p_eye.text = "POWERSCALE ENTERPRISE ARCHITECTURE"
    p_eye.font.size = Pt(14)
    p_eye.font.bold = True
    p_eye.font.color.rgb = BLUE
    p_eye.space_after = Pt(12)

    p_title = tf.add_paragraph()
    p_title.text = "High-Performance Scale-Out Storage Platform"
    p_title.font.size = Pt(40)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_BODY
    p_title.space_after = Pt(16)

    p_sub = tf.add_paragraph()
    p_sub.text = "Unifying file and object workloads with unprecedented scalability, AI data streaming, and zero-trust security."
    p_sub.font.size = Pt(16)
    p_sub.font.color.rgb = TEXT_MUTED

    add_footer(slide1, "01")

    # ==========================================
    # Slide 2: KPI Metrics Dashboard
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2, CANVAS)

    add_header(slide2, "PERFORMANCE METRICS", "PowerScale Next-Gen Cluster Benchmark")

    grid_4col = calculate_grid_positions(cols=4, rows=1, top=Inches(1.8), height=Inches(2.2), gap_x=Inches(0.3))

    kpis = [
        ("186 PB", "Maximum raw capacity per single cluster deployment"),
        ("12.5 TB/s", "Sustained read throughput across NVMe nodes"),
        ("100 GbE", "Dual-port front-end networking interfaces"),
        ("< 1.5 ms", "Ultra-low latency for unstructured AI workloads")
    ]

    for (left, top, w, h), (num, label) in zip(grid_4col, kpis):
        add_kpi_block(slide2, left, top, w, h, num, label)

    # Add a lower 2-column feature block
    grid_2col_bottom = calculate_grid_positions(cols=2, rows=1, top=Inches(4.3), height=Inches(2.4), gap_x=Inches(0.3))

    features = [
        ("OneFS Distributed File System", "Single volume namespace eliminates data silos and management overhead, allowing seamless scaling from terabytes to petabytes."),
        ("Multi-Cloud Integration", "Direct high-speed interconnects with AWS, Azure, and Google Cloud for cloud bursting and automated tiering.")
    ]

    for (left, top, w, h), (title, desc) in zip(grid_2col_bottom, features):
        add_card(slide2, left, top, w, h, bg_color=CARD_BG, border_color=CARD_BORDER)
        tx = slide2.shapes.add_textbox(left + Inches(0.3), top + Inches(0.3), w - Inches(0.6), h - Inches(0.6))
        tf = tx.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0)
        tf.margin_right = Inches(0)
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)

        pt = tf.paragraphs[0]
        pt.text = title
        pt.font.size = Pt(18)
        pt.font.bold = True
        pt.font.color.rgb = TEXT_BODY
        pt.space_after = Pt(8)

        pd = tf.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(13)
        pd.font.color.rgb = TEXT_MUTED

    add_footer(slide2, "02")

    # ==========================================
    # Slide 3: 3-Card Architecture Grid
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3, CANVAS)

    add_header(slide3, "CORE CAPABILITIES", "Architectural Pillars of PowerScale")

    grid_3col = calculate_grid_positions(cols=3, rows=1, top=Inches(1.8), height=Inches(4.9), gap_x=Inches(0.35))

    pillars = [
        ("Flexible All-Flash Nodes", "Powered by NVMe media to accelerate demanding AI, machine learning, and deep learning analytics workloads with massive parallel processing."),
        ("Enterprise Cyber Resilience", "Multi-layer security with SmartLock compliance, automated ransomware detection, air-gap isolation, and immutable snapshots."),
        ("SmartDedupe & Compression", "Inline data reduction technology maximizing storage efficiency and reducing total cost of ownership by up to 33%.")
    ]

    for (left, top, w, h), (p_title, p_desc) in zip(grid_3col, pillars):
        add_card(slide3, left, top, w, h, bg_color=CARD_BG, border_color=CARD_BORDER)

        tx = slide3.shapes.add_textbox(left + Inches(0.3), top + Inches(0.35), w - Inches(0.6), h - Inches(0.7))
        tf = tx.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0)
        tf.margin_right = Inches(0)
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)

        pt = tf.paragraphs[0]
        pt.text = p_title
        pt.font.size = Pt(20)
        pt.font.bold = True
        pt.font.color.rgb = TEXT_BODY
        pt.space_after = Pt(12)

        pd = tf.add_paragraph()
        pd.text = p_desc
        pd.font.size = Pt(14)
        pd.font.color.rgb = TEXT_MUTED

    add_footer(slide3, "03")

    # Save presentation
    prs.save(output_filename)
    print(f"Presentation successfully created and saved to '{output_filename}'")

if __name__ == "__main__":
    build_test_deck()
