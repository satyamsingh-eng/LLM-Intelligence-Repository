"""
PowerScale PPTX Helper Abstractions and Coordinate Math Module
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# 1. Slide Canvas Dimensions (Widescreen 16:9)
SLIDE_WIDTH_INCHES = 13.333
SLIDE_HEIGHT_INCHES = 7.5

SLIDE_WIDTH = Inches(SLIDE_WIDTH_INCHES)
SLIDE_HEIGHT = Inches(SLIDE_HEIGHT_INCHES)

# 2. PowerScale Color Constants
BLUE = RGBColor(0, 102, 255)         # Primary Accent (#0066FF)
BLACK = RGBColor(0, 0, 0)           # Pure Black (#000000)
CANVAS = RGBColor(251, 251, 253)    # Off-white background (#FBFBFD)
CARD_BG = RGBColor(245, 245, 247)   # Light Gray Card Fill (#F5F5F7)
CARD_BORDER = RGBColor(229, 229, 231) # Card Stroke (#E5E5E7)
TEXT_BODY = RGBColor(29, 29, 31)    # Primary Text (#1D1D1F)
TEXT_MUTED = RGBColor(134, 134, 139)# Secondary/Muted Text (#86868B)

# 3. Standard Layout Grids & Margins
MARGIN_LEFT = Inches(0.8)
MARGIN_TOP_HEADER = Inches(0.6)
MARGIN_RIGHT = Inches(0.8)
CONTENT_WIDTH = Inches(11.733)  # 13.333 - 1.6 = 11.733
CONTENT_TOP = Inches(1.8)
CONTENT_HEIGHT = Inches(4.9)    # 7.5 - 1.8 - 0.8 = 4.9
FOOTER_TOP = Inches(7.0)

def create_presentation():
    """Initialize a widescreen 16:9 PowerScale presentation."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs

def set_slide_background(slide, color=CANVAS):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_header(slide, eyebrow_text, title_text, left=MARGIN_LEFT, top=MARGIN_TOP_HEADER, width=CONTENT_WIDTH, height=Inches(1.0)):
    """
    Add a standard PowerScale header block with eyebrow and main title.
    Calculated to avoid text clipping or line wrapping issues.
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)

    # Eyebrow
    p_eyebrow = tf.paragraphs[0]
    p_eyebrow.text = eyebrow_text.upper() if eyebrow_text else ""
    p_eyebrow.font.size = Pt(12)
    p_eyebrow.font.bold = True
    p_eyebrow.font.color.rgb = BLUE
    p_eyebrow.font.name = "Calibri"
    p_eyebrow.space_after = Pt(4)

    # Title
    p_title = tf.add_paragraph()
    p_title.text = title_text
    p_title.font.size = Pt(32)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_BODY
    p_title.font.name = "Calibri"

    return txBox

def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=CARD_BORDER):
    """
    Add a card container rectangle with solid background and stroke.
    """
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color

    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()

    return shape

def add_kpi_block(slide, left, top, width, height, number_str, label_text, bg_color=CARD_BG, border_color=CARD_BORDER, num_color=BLUE):
    """
    Add a KPI block card container with big number stat and label text.
    Uses precise internal padding and vertical spacing.
    """
    card = add_card(slide, left, top, width, height, bg_color, border_color)

    pad_x = Inches(0.25)
    pad_y = Inches(0.25)

    txBox = slide.shapes.add_textbox(
        left + pad_x,
        top + pad_y,
        width - (2 * pad_x),
        height - (2 * pad_y)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)

    # Big Number
    p_num = tf.paragraphs[0]
    p_num.text = number_str
    p_num.font.size = Pt(44)
    p_num.font.bold = True
    p_num.font.color.rgb = num_color
    p_num.font.name = "Calibri"
    p_num.space_after = Pt(6)

    # Label
    p_label = tf.add_paragraph()
    p_label.text = label_text
    p_label.font.size = Pt(13)
    p_label.font.bold = False
    p_label.font.color.rgb = TEXT_BODY
    p_label.font.name = "Calibri"

    return card

def add_footer(slide, slide_num_str, left=MARGIN_LEFT, top=FOOTER_TOP, width=CONTENT_WIDTH, height=Inches(0.3)):
    """
    Add a standardized footer block containing brand note and slide number.
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)

    p = tf.paragraphs[0]
    p.font.size = Pt(10)
    p.font.name = "Calibri"
    p.font.color.rgb = TEXT_MUTED

    run1 = p.add_run()
    run1.text = "PowerScale Platform Architecture"

    # Add space and right-aligned slide number via second text box
    txBox_num = slide.shapes.add_textbox(left + width - Inches(2.0), top, Inches(2.0), height)
    tf_num = txBox_num.text_frame
    tf_num.word_wrap = False
    tf_num.margin_left = Inches(0)
    tf_num.margin_right = Inches(0)
    tf_num.margin_top = Inches(0)
    tf_num.margin_bottom = Inches(0)

    p_num = tf_num.paragraphs[0]
    p_num.alignment = PP_ALIGN.RIGHT
    p_num.font.size = Pt(10)
    p_num.font.name = "Calibri"
    p_num.font.color.rgb = TEXT_MUTED
    p_num.text = str(slide_num_str)

    return txBox

def calculate_grid_positions(cols, rows, left=MARGIN_LEFT, top=CONTENT_TOP, width=CONTENT_WIDTH, height=CONTENT_HEIGHT, gap_x=Inches(0.3), gap_y=Inches(0.3)):
    """
    Calculate card bounds (left, top, card_w, card_h) for a grid layout.
    Guarantees no overlap and uniform margins across columns and rows.
    """
    total_gap_x = gap_x * (cols - 1)
    card_w = (width - total_gap_x) / cols

    total_gap_y = gap_y * (rows - 1)
    card_h = (height - total_gap_y) / rows

    coords = []
    for r in range(rows):
        for c in range(cols):
            c_left = left + c * (card_w + gap_x)
            c_top = top + r * (card_h + gap_y)
            coords.append((c_left, c_top, card_w, card_h))

    return coords
